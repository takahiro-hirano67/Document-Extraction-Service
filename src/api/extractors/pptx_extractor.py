# src/api/extractors/pptx_extractor.py

"""PowerPointファイル(.pptx)からテキストを抽出するモジュール

スライド内のテキスト（テキストボックス、図形、表）およびスライドノートを抽出します。
グループ化された図形内のテキストも再帰的に取得します。
"""

from __future__ import annotations

import re
from typing import TYPE_CHECKING, Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

if TYPE_CHECKING:
    import io
    from collections.abc import Iterator

    from pptx.shapes.base import BaseShape
    from pptx.table import Table


def extract_from_pptx(file_stream: io.BytesIO) -> str:
    """PowerPointファイルからテキストを抽出する

    抽出順序:
        1. 各スライド内の図形・テキストボックス・表（グループ化図形含む）
        2. スライドノート（存在する場合）

    Args:
        file_stream (io.BytesIO): PowerPointファイルのバイトストリーム

    Returns:
        str: 抽出されたテキスト全文

    """
    presentation = Presentation(file_stream)
    extracted_lines: list[str] = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        extracted_lines.append(f"\n=== スライド {slide_index} ===\n")

        # 図形のY座標・X座標を取得するヘルパー関数
        def _get_shape_position(shape: BaseShape) -> tuple[int, int]:
            top = getattr(shape, "top", 0) or 0
            left = getattr(shape, "left", 0) or 0
            # 多少の縦のズレを吸収し、同じ行とみなすためにY座標を丸める (100,000 EMU ≒ 約2.8mm)
            return (round(top / 100000), left)

        # スライド内の図形を「上から下」「左から右」の順にソートする
        sorted_shapes = sorted(slide.shapes, key=_get_shape_position)

        # ソート済みの図形を順次走査
        for shape in sorted_shapes:
            extracted_lines.extend(_extract_text_from_shape(shape))

        # スライドノートの抽出
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                extracted_lines.append("\n【ノート】")
                extracted_lines.append(notes_text)

    # 各行を改行で結合し、連続する空白行を1行に正規化する
    joined_text = "\n".join(extracted_lines)
    normalized_text = re.sub(r"\n{3,}", "\n\n", joined_text)

    return normalized_text.strip()


def _extract_text_from_shape(shape: BaseShape) -> Iterator[str]:
    """図形要素からテキストを抽出する（グループ図形の場合は再帰処理）

    Pylance対策として、BaseShapeに定義されていない属性(shapes, table, text等)
    へのアクセスには getattr() を使用して型エラーを回避します。

    Args:
        shape (BaseShape): python-pptxの図形オブジェクト

    Yields:
        str: 抽出されたテキスト行

    """
    # 1. グループ化された図形の場合、内部の図形を再帰的に処理
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        child_shapes = getattr(shape, "shapes", [])
        for child_shape in child_shapes:
            yield from _extract_text_from_shape(child_shape)

    # 2. 表の場合、Markdown形式の表に変換
    elif getattr(shape, "has_table", False):
        table = getattr(shape, "table", None)
        if table is not None:
            markdown_table = _table_to_markdown(table)
            if markdown_table:
                yield f"\n{markdown_table}\n"

    # 3. 通常のテキストフレーム（テキストボックス、図形内テキスト）の場合
    elif getattr(shape, "has_text_frame", False):
        text = getattr(shape, "text", "")
        if text:
            yield text.strip()


def _table_to_markdown(table: Table) -> str:
    """python-pptxのTableオブジェクトをMarkdownの表形式に変換する

    Args:
        table (Table): python-pptxのTableオブジェクト

    Returns:
        str: Markdown形式の表文字列。行が存在しない場合は空文字列を返す。

    """
    # Pylanceで _RowCollection が Iterable と認識されないため Any で受ける
    rows: Any = getattr(table, "rows", [])
    if not rows:
        return ""

    markdown_lines: list[str] = []

    for row_index, row in enumerate(rows):
        row_cells: list[str] = []
        cells: Any = getattr(row, "cells", [])

        for cell in cells:
            text_frame = getattr(cell, "text_frame", None)
            cell_text = getattr(text_frame, "text", "") if text_frame else ""
            row_cells.append(cell_text.replace("\n", " ").replace("|", "\\|").strip())

        markdown_lines.append("| " + " | ".join(row_cells) + " |")

        # 1行目（ヘッダー）の直後にMarkdownの区切り線を追加
        if row_index == 0:
            separator_cells = ["---"] * len(row_cells)
            markdown_lines.append("| " + " | ".join(separator_cells) + " |")

    return "\n".join(markdown_lines)
